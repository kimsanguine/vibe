"""
Export Skill
- python-pptx를 사용한 PPTX 파일 생성
- 한국어 폰트 지원
- 기본 레이아웃 및 스타일 적용
"""

import os
from typing import Optional
from pathlib import Path

from ..config import AgentConfig, get_config
from ..models.slide import Presentation, Slide, SlideType, DesignSystem


class ExportSkill:
    """
    Export Skill

    역할:
    1. python-pptx를 사용하여 PPTX 파일 생성
    2. 슬라이드 타입별 레이아웃 적용
    3. 디자인 시스템(색상, 폰트) 적용
    """

    def __init__(self, config: Optional[AgentConfig] = None):
        self.config = config or get_config()
        self._pptx = None
        self._init_pptx()

    def _init_pptx(self):
        """python-pptx 초기화"""
        try:
            from pptx import Presentation as PptxPresentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
            self._pptx = PptxPresentation
            self._Inches = Inches
            self._Pt = Pt
            self._RGBColor = RGBColor
            self._PP_ALIGN = PP_ALIGN
            self._MSO_ANCHOR = MSO_ANCHOR
        except ImportError:
            raise ImportError("python-pptx 패키지가 필요합니다: pip install python-pptx")

    def _hex_to_rgb(self, hex_color: str) -> tuple:
        """HEX 색상을 RGB로 변환"""
        hex_color = hex_color.lstrip('#')
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

    def _apply_design_to_shape(self, shape, design: DesignSystem, is_title: bool = False):
        """Shape에 디자인 적용"""
        from pptx.dml.color import RGBColor

        if hasattr(shape, "text_frame"):
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = design.font_title if is_title else design.font_body
                    run.font.size = self._Pt(design.font_size_title if is_title else design.font_size_body)

                    # 텍스트 색상
                    rgb = self._hex_to_rgb(design.text_color)
                    run.font.color.rgb = RGBColor(*rgb)

    def _create_title_slide(self, prs, slide: Slide, design: DesignSystem):
        """타이틀 슬라이드 생성"""
        slide_layout = prs.slide_layouts[6]  # 빈 레이아웃
        pptx_slide = prs.slides.add_slide(slide_layout)

        # 배경색 설정
        from pptx.dml.color import RGBColor
        background = pptx_slide.background
        fill = background.fill
        fill.solid()
        rgb = self._hex_to_rgb(design.primary_color)
        fill.fore_color.rgb = RGBColor(*rgb)

        # 제목
        left = self._Inches(1)
        top = self._Inches(3)
        width = self._Inches(8)
        height = self._Inches(1.5)

        title_box = pptx_slide.shapes.add_textbox(left, top, width, height)
        tf = title_box.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = slide.content.title
        p.font.name = design.font_title
        p.font.size = self._Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)  # 흰색
        p.alignment = self._PP_ALIGN.CENTER

        # 부제목 (body가 있는 경우)
        if slide.content.body:
            subtitle_top = self._Inches(4.5)
            subtitle_box = pptx_slide.shapes.add_textbox(left, subtitle_top, width, self._Inches(0.5))
            stf = subtitle_box.text_frame
            sp = stf.paragraphs[0]
            sp.text = slide.content.body[0]
            sp.font.name = design.font_body
            sp.font.size = self._Pt(24)
            sp.font.color.rgb = RGBColor(255, 255, 255)
            sp.alignment = self._PP_ALIGN.CENTER

        return pptx_slide

    def _create_content_slide(self, prs, slide: Slide, design: DesignSystem):
        """콘텐츠 슬라이드 생성"""
        from pptx.dml.color import RGBColor

        slide_layout = prs.slide_layouts[6]  # 빈 레이아웃
        pptx_slide = prs.slides.add_slide(slide_layout)

        # 제목
        left = self._Inches(0.5)
        top = self._Inches(0.5)
        width = self._Inches(9)
        height = self._Inches(1)

        title_box = pptx_slide.shapes.add_textbox(left, top, width, height)
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = slide.content.title
        p.font.name = design.font_title
        p.font.size = self._Pt(32)
        p.font.bold = True
        rgb = self._hex_to_rgb(design.primary_color)
        p.font.color.rgb = RGBColor(*rgb)

        # 본문 내용
        content_top = self._Inches(1.8)
        content_height = self._Inches(5)

        content_box = pptx_slide.shapes.add_textbox(left, content_top, width, content_height)
        ctf = content_box.text_frame
        ctf.word_wrap = True

        for i, bullet in enumerate(slide.content.body):
            if i == 0:
                cp = ctf.paragraphs[0]
            else:
                cp = ctf.add_paragraph()

            cp.text = f"• {bullet}"
            cp.font.name = design.font_body
            cp.font.size = self._Pt(18)
            cp.space_after = self._Pt(12)
            rgb = self._hex_to_rgb(design.text_color)
            cp.font.color.rgb = RGBColor(*rgb)

        # 노트 추가
        if slide.content.notes:
            notes_slide = pptx_slide.notes_slide
            notes_slide.notes_text_frame.text = slide.content.notes

        return pptx_slide

    def _create_section_slide(self, prs, slide: Slide, design: DesignSystem):
        """섹션 구분 슬라이드 생성"""
        from pptx.dml.color import RGBColor

        slide_layout = prs.slide_layouts[6]
        pptx_slide = prs.slides.add_slide(slide_layout)

        # 배경색
        background = pptx_slide.background
        fill = background.fill
        fill.solid()
        rgb = self._hex_to_rgb(design.secondary_color)
        fill.fore_color.rgb = RGBColor(*rgb)

        # 섹션 제목
        left = self._Inches(1)
        top = self._Inches(3)
        width = self._Inches(8)
        height = self._Inches(1.5)

        title_box = pptx_slide.shapes.add_textbox(left, top, width, height)
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = slide.content.title
        p.font.name = design.font_title
        p.font.size = self._Pt(40)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = self._PP_ALIGN.CENTER

        return pptx_slide

    def _create_conclusion_slide(self, prs, slide: Slide, design: DesignSystem):
        """결론 슬라이드 생성"""
        from pptx.dml.color import RGBColor

        slide_layout = prs.slide_layouts[6]
        pptx_slide = prs.slides.add_slide(slide_layout)

        # 제목
        left = self._Inches(0.5)
        top = self._Inches(0.5)
        width = self._Inches(9)
        height = self._Inches(1)

        title_box = pptx_slide.shapes.add_textbox(left, top, width, height)
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = slide.content.title
        p.font.name = design.font_title
        p.font.size = self._Pt(32)
        p.font.bold = True
        rgb = self._hex_to_rgb(design.primary_color)
        p.font.color.rgb = RGBColor(*rgb)

        # 핵심 요약
        content_top = self._Inches(2)
        content_box = pptx_slide.shapes.add_textbox(left, content_top, width, self._Inches(4))
        ctf = content_box.text_frame
        ctf.word_wrap = True

        for i, point in enumerate(slide.content.body):
            if i == 0:
                cp = ctf.paragraphs[0]
            else:
                cp = ctf.add_paragraph()

            cp.text = f"✓ {point}"
            cp.font.name = design.font_body
            cp.font.size = self._Pt(20)
            cp.font.bold = True
            cp.space_after = self._Pt(16)
            rgb = self._hex_to_rgb(design.text_color)
            cp.font.color.rgb = RGBColor(*rgb)

        return pptx_slide

    def export(
        self,
        presentation: Presentation,
        output_path: Optional[str] = None
    ) -> str:
        """PPTX 파일로 내보내기"""

        # 출력 경로 설정
        if output_path is None:
            output_dir = Path(self.config.output_dir)
            output_dir.mkdir(parents=True, exist_ok=True)

            # 파일명 생성 (주제에서 안전한 문자만 추출)
            safe_topic = "".join(c for c in presentation.topic[:30] if c.isalnum() or c in " _-")
            safe_topic = safe_topic.strip().replace(" ", "_")
            output_path = str(output_dir / f"{safe_topic}.pptx")

        # PPTX 생성
        prs = self._pptx()

        # 슬라이드 크기 설정 (16:9)
        prs.slide_width = self._Inches(10)
        prs.slide_height = self._Inches(5.625)

        # 각 슬라이드 생성
        for slide in presentation.slides:
            if slide.slide_type == SlideType.TITLE:
                self._create_title_slide(prs, slide, presentation.design)
            elif slide.slide_type == SlideType.SECTION:
                self._create_section_slide(prs, slide, presentation.design)
            elif slide.slide_type == SlideType.CONCLUSION:
                self._create_conclusion_slide(prs, slide, presentation.design)
            else:
                self._create_content_slide(prs, slide, presentation.design)

        # 저장
        prs.save(output_path)

        print(f"\n[Export Skill] PPTX 파일 생성 완료: {output_path}")
        print(f"[슬라이드 수] {len(presentation.slides)}장")

        return output_path
